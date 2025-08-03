import asyncio
import json
import logging
from datetime import datetime
from typing import Dict, List, Optional, Any
from pathlib import Path
import pandas as pd
import numpy as np

# AI/ML imports
from transformers import pipeline, AutoTokenizer, AutoModelForTextGeneration
import openai

# Presentation generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Data visualization
import matplotlib.pyplot as plt
import seaborn as sns
from io import BytesIO

# Custom imports
from ...backend.app.models.deal import Deal, DealType
from ...backend.app.models.document import Document
from ..data_processing.refinitiv_api.refinitiv_client import RefinitivClient
from ..data_processing.sec_edgar.sec_client import SECClient

logger = logging.getLogger(__name__)


class PitchbookGenerator:
    """AI-powered pitchbook generator for investment banking deals"""
    
    def __init__(self):
        self.refinitiv_client = RefinitivClient()
        self.sec_client = SECClient()
        self.text_generator = None
        self.summarizer = None
        self._load_models()
    
    def _load_models(self):
        """Load AI models for text generation and summarization"""
        try:
            # Load text generation model
            self.text_generator = pipeline(
                "text-generation",
                model="gpt2",  # In production, use a more sophisticated model
                tokenizer="gpt2"
            )
            
            # Load summarization model
            self.summarizer = pipeline(
                "summarization",
                model="facebook/bart-large-cnn"
            )
            
            logger.info("Pitchbook generation models loaded successfully")
            
        except Exception as e:
            logger.error(f"Error loading pitchbook generation models: {e}")
            raise
    
    async def generate_pitchbook(self, deal: Deal, documents: List[Document] = None) -> Dict[str, Any]:
        """Generate a complete pitchbook for a deal"""
        start_time = datetime.utcnow()
        
        try:
            logger.info(f"Starting pitchbook generation for deal {deal.id}")
            
            # Gather market data
            market_data = await self._gather_market_data(deal)
            
            # Generate slide content
            slides_content = await self._generate_slides_content(deal, market_data, documents)
            
            # Create PowerPoint presentation
            presentation_path = await self._create_presentation(deal, slides_content)
            
            # Generate executive summary
            executive_summary = await self._generate_executive_summary(deal, slides_content)
            
            processing_time = (datetime.utcnow() - start_time).total_seconds()
            logger.info(f"Pitchbook generation completed in {processing_time:.2f} seconds")
            
            return {
                'success': True,
                'processing_time': processing_time,
                'presentation_path': presentation_path,
                'executive_summary': executive_summary,
                'slides_count': len(slides_content),
                'market_data_used': len(market_data)
            }
            
        except Exception as e:
            logger.error(f"Error in pitchbook generation: {e}")
            return {
                'success': False,
                'error': str(e),
                'processing_time': (datetime.utcnow() - start_time).total_seconds()
            }
    
    async def _gather_market_data(self, deal: Deal) -> Dict[str, Any]:
        """Gather relevant market data for the deal"""
        try:
            market_data = {}
            
            # Get industry data
            if deal.target_industry:
                industry_data = await self.refinitiv_client.get_industry_data(deal.target_industry)
                market_data['industry'] = industry_data
            
            # Get comparable companies
            if deal.target_company:
                comps_data = await self.refinitiv_client.get_comparable_companies(deal.target_company)
                market_data['comparable_companies'] = comps_data
            
            # Get market trends
            market_trends = await self.refinitiv_client.get_market_trends(deal.target_sector)
            market_data['market_trends'] = market_trends
            
            # Get SEC filings if available
            if deal.target_company:
                sec_filings = await self.sec_client.get_company_filings(deal.target_company)
                market_data['sec_filings'] = sec_filings
            
            return market_data
            
        except Exception as e:
            logger.error(f"Error gathering market data: {e}")
            return {}
    
    async def _generate_slides_content(self, deal: Deal, market_data: Dict[str, Any], 
                                     documents: List[Document] = None) -> List[Dict[str, Any]]:
        """Generate content for each slide"""
        try:
            slides = []
            
            # Slide 1: Title Slide
            slides.append(await self._generate_title_slide(deal))
            
            # Slide 2: Executive Summary
            slides.append(await self._generate_executive_summary_slide(deal))
            
            # Slide 3: Deal Overview
            slides.append(await self._generate_deal_overview_slide(deal))
            
            # Slide 4: Company Overview
            if deal.target_company:
                slides.append(await self._generate_company_overview_slide(deal, market_data))
            
            # Slide 5: Financial Highlights
            slides.append(await self._generate_financial_highlights_slide(deal, market_data))
            
            # Slide 6: Market Analysis
            slides.append(await self._generate_market_analysis_slide(deal, market_data))
            
            # Slide 7: Comparable Companies
            if market_data.get('comparable_companies'):
                slides.append(await self._generate_comparable_companies_slide(market_data))
            
            # Slide 8: Transaction Structure
            slides.append(await self._generate_transaction_structure_slide(deal))
            
            # Slide 9: Valuation Analysis
            slides.append(await self._generate_valuation_analysis_slide(deal, market_data))
            
            # Slide 10: Risk Analysis
            slides.append(await self._generate_risk_analysis_slide(deal, documents))
            
            # Slide 11: Timeline
            slides.append(await self._generate_timeline_slide(deal))
            
            # Slide 12: Next Steps
            slides.append(await self._generate_next_steps_slide(deal))
            
            return slides
            
        except Exception as e:
            logger.error(f"Error generating slides content: {e}")
            return []
    
    async def _generate_title_slide(self, deal: Deal) -> Dict[str, Any]:
        """Generate title slide content"""
        return {
            'slide_number': 1,
            'title': f"{deal.name} - Investment Opportunity",
            'subtitle': f"{deal.deal_type.value.upper()} Transaction",
            'type': 'title',
            'content': {
                'company_name': deal.target_company or 'Target Company',
                'deal_type': deal.deal_type.value.upper(),
                'date': datetime.utcnow().strftime('%B %Y')
            }
        }
    
    async def _generate_executive_summary_slide(self, deal: Deal) -> Dict[str, Any]:
        """Generate executive summary slide"""
        summary_points = [
            f"Deal Value: ${deal.deal_value:.1f}M" if deal.deal_value else "Deal Value: TBD",
            f"Target Company: {deal.target_company}" if deal.target_company else "Target: Confidential",
            f"Industry: {deal.target_industry}" if deal.target_industry else "Industry: TBD",
            f"Expected Close: {deal.expected_close_date.strftime('%Q4 %Y')}" if deal.expected_close_date else "Timeline: TBD"
        ]
        
        return {
            'slide_number': 2,
            'title': 'Executive Summary',
            'type': 'bullet_points',
            'content': {
                'points': summary_points,
                'key_highlights': [
                    'Strategic investment opportunity',
                    'Strong market positioning',
                    'Attractive valuation metrics',
                    'Clear path to value creation'
                ]
            }
        }
    
    async def _generate_deal_overview_slide(self, deal: Deal) -> Dict[str, Any]:
        """Generate deal overview slide"""
        return {
            'slide_number': 3,
            'title': 'Deal Overview',
            'type': 'overview',
            'content': {
                'deal_type': deal.deal_type.value,
                'target_company': deal.target_company,
                'target_industry': deal.target_industry,
                'target_sector': deal.target_sector,
                'deal_value': deal.deal_value,
                'currency': deal.deal_currency,
                'transaction_fee': deal.transaction_fee,
                'success_fee_rate': deal.success_fee_rate
            }
        }
    
    async def _generate_company_overview_slide(self, deal: Deal, market_data: Dict[str, Any]) -> Dict[str, Any]:
        """Generate company overview slide"""
        company_info = market_data.get('industry', {})
        
        return {
            'slide_number': 4,
            'title': 'Company Overview',
            'type': 'company_profile',
            'content': {
                'company_name': deal.target_company,
                'industry': deal.target_industry,
                'sector': deal.target_sector,
                'revenue': deal.target_revenue,
                'ebitda': deal.target_ebitda,
                'description': company_info.get('description', 'Company description to be provided'),
                'key_metrics': company_info.get('key_metrics', [])
            }
        }
    
    async def _generate_financial_highlights_slide(self, deal: Deal, market_data: Dict[str, Any]) -> Dict[str, Any]:
        """Generate financial highlights slide"""
        financial_data = []
        
        if deal.target_revenue:
            financial_data.append({
                'metric': 'Revenue',
                'value': f"${deal.target_revenue:.1f}M",
                'period': 'LTM'
            })
        
        if deal.target_ebitda:
            financial_data.append({
                'metric': 'EBITDA',
                'value': f"${deal.target_ebitda:.1f}M",
                'period': 'LTM'
            })
        
        return {
            'slide_number': 5,
            'title': 'Financial Highlights',
            'type': 'financial_metrics',
            'content': {
                'metrics': financial_data,
                'charts': ['revenue_trend', 'ebitda_margin'],
                'key_insights': [
                    'Strong revenue growth trajectory',
                    'Improving profitability margins',
                    'Healthy cash flow generation'
                ]
            }
        }
    
    async def _generate_market_analysis_slide(self, deal: Deal, market_data: Dict[str, Any]) -> Dict[str, Any]:
        """Generate market analysis slide"""
        market_trends = market_data.get('market_trends', {})
        
        return {
            'slide_number': 6,
            'title': 'Market Analysis',
            'type': 'market_analysis',
            'content': {
                'market_size': market_trends.get('market_size'),
                'growth_rate': market_trends.get('growth_rate'),
                'key_drivers': market_trends.get('key_drivers', []),
                'competitive_landscape': market_trends.get('competitive_landscape', []),
                'regulatory_environment': market_trends.get('regulatory_environment', [])
            }
        }
    
    async def _generate_comparable_companies_slide(self, market_data: Dict[str, Any]) -> Dict[str, Any]:
        """Generate comparable companies slide"""
        comps_data = market_data.get('comparable_companies', [])
        
        return {
            'slide_number': 7,
            'title': 'Comparable Companies',
            'type': 'comparable_companies',
            'content': {
                'companies': comps_data[:5],  # Top 5 comparables
                'valuation_metrics': ['EV/Revenue', 'EV/EBITDA', 'P/E'],
                'average_metrics': self._calculate_average_metrics(comps_data)
            }
        }
    
    async def _generate_transaction_structure_slide(self, deal: Deal) -> Dict[str, Any]:
        """Generate transaction structure slide"""
        return {
            'slide_number': 8,
            'title': 'Transaction Structure',
            'type': 'transaction_structure',
            'content': {
                'deal_type': deal.deal_type.value,
                'structure': self._get_deal_structure(deal.deal_type),
                'consideration': self._get_consideration_structure(deal),
                'financing': 'To be determined',
                'closing_conditions': [
                    'Regulatory approvals',
                    'Due diligence completion',
                    'Financing arrangements',
                    'Shareholder approval'
                ]
            }
        }
    
    async def _generate_valuation_analysis_slide(self, deal: Deal, market_data: Dict[str, Any]) -> Dict[str, Any]:
        """Generate valuation analysis slide"""
        comps_data = market_data.get('comparable_companies', [])
        
        return {
            'slide_number': 9,
            'title': 'Valuation Analysis',
            'type': 'valuation',
            'content': {
                'valuation_methods': ['Comparable Companies', 'Precedent Transactions', 'DCF'],
                'implied_valuation': deal.deal_value,
                'valuation_range': self._calculate_valuation_range(deal, comps_data),
                'key_assumptions': [
                    'Revenue growth assumptions',
                    'Margin expansion potential',
                    'Exit multiple scenarios'
                ]
            }
        }
    
    async def _generate_risk_analysis_slide(self, deal: Deal, documents: List[Document] = None) -> Dict[str, Any]:
        """Generate risk analysis slide"""
        risk_factors = [
            'Market and economic conditions',
            'Regulatory changes',
            'Competitive pressures',
            'Integration challenges',
            'Financing risks'
        ]
        
        return {
            'slide_number': 10,
            'title': 'Risk Analysis',
            'type': 'risk_analysis',
            'content': {
                'risk_factors': risk_factors,
                'mitigation_strategies': [
                    'Comprehensive due diligence',
                    'Structured transaction terms',
                    'Experienced management team',
                    'Robust integration plan'
                ],
                'risk_score': 'Medium to Low'
            }
        }
    
    async def _generate_timeline_slide(self, deal: Deal) -> Dict[str, Any]:
        """Generate timeline slide"""
        timeline = [
            {'phase': 'Due Diligence', 'duration': '4-6 weeks', 'status': 'In Progress'},
            {'phase': 'Negotiation', 'duration': '2-4 weeks', 'status': 'Pending'},
            {'phase': 'Documentation', 'duration': '2-3 weeks', 'status': 'Pending'},
            {'phase': 'Regulatory Approval', 'duration': '4-8 weeks', 'status': 'Pending'},
            {'phase': 'Closing', 'duration': '1 week', 'status': 'Pending'}
        ]
        
        return {
            'slide_number': 11,
            'title': 'Transaction Timeline',
            'type': 'timeline',
            'content': {
                'timeline': timeline,
                'expected_close': deal.expected_close_date,
                'key_milestones': [
                    'Due diligence completion',
                    'Definitive agreement signing',
                    'Regulatory approvals',
                    'Closing'
                ]
            }
        }
    
    async def _generate_next_steps_slide(self, deal: Deal) -> Dict[str, Any]:
        """Generate next steps slide"""
        return {
            'slide_number': 12,
            'title': 'Next Steps',
            'type': 'next_steps',
            'content': {
                'immediate_actions': [
                    'Complete due diligence review',
                    'Finalize transaction structure',
                    'Engage legal counsel',
                    'Prepare definitive agreements'
                ],
                'timeline': 'Next 30 days',
                'key_contacts': [
                    'Investment Banking Team',
                    'Legal Counsel',
                    'Financial Advisors'
                ]
            }
        }
    
    async def _create_presentation(self, deal: Deal, slides_content: List[Dict[str, Any]]) -> str:
        """Create PowerPoint presentation from slide content"""
        try:
            # Create presentation
            prs = Presentation()
            
            # Set slide size to 16:9
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            for slide_content in slides_content:
                await self._add_slide_to_presentation(prs, slide_content)
            
            # Save presentation
            filename = f"pitchbook_{deal.id}_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.pptx"
            filepath = f"/tmp/{filename}"
            prs.save(filepath)
            
            return filepath
            
        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            raise
    
    async def _add_slide_to_presentation(self, prs: Presentation, slide_content: Dict[str, Any]):
        """Add a slide to the presentation"""
        try:
            slide_type = slide_content.get('type', 'content')
            
            if slide_type == 'title':
                slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                title.text = slide_content['title']
                subtitle.text = slide_content['content']['company_name']
                
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # Content slide
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = slide_content['title']
                
                # Add content based on slide type
                if slide_type == 'bullet_points':
                    content_text = '\n'.join([f"• {point}" for point in slide_content['content']['points']])
                    content.text = content_text
                
                elif slide_type == 'financial_metrics':
                    metrics_text = '\n'.join([
                        f"• {metric['metric']}: {metric['value']}"
                        for metric in slide_content['content']['metrics']
                    ])
                    content.text = metrics_text
                
                else:
                    content.text = str(slide_content['content'])
            
            # Apply formatting
            self._format_slide(slide)
            
        except Exception as e:
            logger.error(f"Error adding slide: {e}")
    
    def _format_slide(self, slide):
        """Apply formatting to slide"""
        try:
            # Format title
            if slide.shapes.title:
                title = slide.shapes.title
                title.text_frame.paragraphs[0].font.size = Pt(44)
                title.text_frame.paragraphs[0].font.bold = True
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            # Format content
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(18)
                        paragraph.font.color.rgb = RGBColor(51, 51, 51)
            
        except Exception as e:
            logger.error(f"Error formatting slide: {e}")
    
    async def _generate_executive_summary(self, deal: Deal, slides_content: List[Dict[str, Any]]) -> str:
        """Generate executive summary text"""
        try:
            summary = f"Executive Summary: {deal.name}\n\n"
            summary += f"Transaction Type: {deal.deal_type.value.upper()}\n"
            
            if deal.deal_value:
                summary += f"Deal Value: ${deal.deal_value:.1f}M\n"
            
            if deal.target_company:
                summary += f"Target Company: {deal.target_company}\n"
            
            summary += f"Industry: {deal.target_industry or 'TBD'}\n"
            summary += f"Expected Close: {deal.expected_close_date.strftime('%Q4 %Y') if deal.expected_close_date else 'TBD'}\n\n"
            
            summary += "Key Highlights:\n"
            summary += "• Strategic investment opportunity with strong growth potential\n"
            summary += "• Attractive valuation metrics relative to peers\n"
            summary += "• Clear path to value creation through operational improvements\n"
            summary += "• Experienced management team with proven track record\n"
            
            return summary
            
        except Exception as e:
            logger.error(f"Error generating executive summary: {e}")
            return "Executive summary generation failed"
    
    def _calculate_average_metrics(self, comps_data: List[Dict[str, Any]]) -> Dict[str, float]:
        """Calculate average valuation metrics for comparable companies"""
        try:
            if not comps_data:
                return {}
            
            metrics = {}
            for comp in comps_data:
                for key, value in comp.get('valuation_metrics', {}).items():
                    if key not in metrics:
                        metrics[key] = []
                    if value and value > 0:
                        metrics[key].append(value)
            
            return {key: np.mean(values) for key, values in metrics.items() if values}
            
        except Exception as e:
            logger.error(f"Error calculating average metrics: {e}")
            return {}
    
    def _get_deal_structure(self, deal_type: DealType) -> str:
        """Get deal structure based on deal type"""
        structures = {
            DealType.MNA: "Stock Purchase / Asset Purchase",
            DealType.IPO: "Initial Public Offering",
            DealType.PRIVATE_EQUITY: "Leveraged Buyout / Growth Investment",
            DealType.DEBT_FINANCING: "Senior / Subordinated Debt",
            DealType.RESTRUCTURING: "Debt Restructuring / Equity Conversion",
            DealType.OTHER: "Custom Structure"
        }
        return structures.get(deal_type, "To be determined")
    
    def _get_consideration_structure(self, deal: Deal) -> str:
        """Get consideration structure for the deal"""
        if deal.deal_type == DealType.MNA:
            return "Cash / Stock / Mixed"
        elif deal.deal_type == DealType.IPO:
            return "Primary / Secondary Shares"
        elif deal.deal_type == DealType.PRIVATE_EQUITY:
            return "Equity / Debt / Preferred"
        else:
            return "To be determined"
    
    def _calculate_valuation_range(self, deal: Deal, comps_data: List[Dict[str, Any]]) -> Dict[str, float]:
        """Calculate valuation range based on comparable companies"""
        try:
            if not comps_data or not deal.target_revenue:
                return {}
            
            # Calculate EV/Revenue multiples
            ev_revenue_multiples = []
            for comp in comps_data:
                if comp.get('ev_revenue_multiple') and comp['ev_revenue_multiple'] > 0:
                    ev_revenue_multiples.append(comp['ev_revenue_multiple'])
            
            if ev_revenue_multiples:
                avg_multiple = np.mean(ev_revenue_multiples)
                std_multiple = np.std(ev_revenue_multiples)
                
                low_multiple = avg_multiple - std_multiple
                high_multiple = avg_multiple + std_multiple
                
                return {
                    'low_valuation': deal.target_revenue * low_multiple,
                    'high_valuation': deal.target_revenue * high_multiple,
                    'average_multiple': avg_multiple
                }
            
            return {}
            
        except Exception as e:
            logger.error(f"Error calculating valuation range: {e}")
            return {} 